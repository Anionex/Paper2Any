from __future__ import annotations

import asyncio
import json
import time
from pathlib import Path
from typing import List, Dict, Any
from urllib.parse import urlparse

import fitz

from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger
from dataflow_agent.state import KBReportState, MainState
from dataflow_agent.agentroles import create_agent
from dataflow_agent.utils import get_project_root

log = get_logger(__name__)

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def _resolve_local_path(path_or_url: str) -> Path:
    if not path_or_url:
        return Path()
    if path_or_url.startswith("http"):
        parsed = urlparse(path_or_url)
        path_or_url = parsed.path
    if path_or_url.startswith("/outputs/"):
        rel = path_or_url[len("/outputs/") :].lstrip("/")
        return (get_project_root() / "outputs" / rel).resolve()
    p = Path(path_or_url)
    if p.is_absolute():
        return p.resolve()
    return (get_project_root() / p).resolve()


def _ensure_under_outputs(path: Path) -> bool:
    outputs_root = (get_project_root() / "outputs").resolve()
    try:
        path.resolve().relative_to(outputs_root)
        return True
    except Exception:
        return False


def _extract_file_text(path: Path) -> str:
    if not path.exists():
        return ""
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            doc = fitz.open(path)
            text = ""
            for page in doc:
                text += page.get_text() + "\n"
            return text
        if suffix in {".docx", ".doc"}:
            if Document is None:
                return ""
            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs])
        if suffix in {".pptx", ".ppt"}:
            if Presentation is None:
                return ""
            prs = Presentation(path)
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        if suffix in {".md", ".txt"}:
            return path.read_text(encoding="utf-8", errors="ignore")
        return ""
    except Exception:
        return ""


def _build_file_entries(file_paths: List[str], per_file_chars: int = 20000) -> List[Dict[str, Any]]:
    entries: List[Dict[str, Any]] = []
    for raw in file_paths:
        try:
            local_path = _resolve_local_path(raw)
            if not _ensure_under_outputs(local_path):
                continue
            text = _extract_file_text(local_path)
            if not text:
                continue
            if per_file_chars and len(text) > per_file_chars:
                text = text[:per_file_chars] + "\n\n[...content truncated]"
            entries.append({
                "file_name": local_path.name,
                "path": str(local_path),
                "content": text
            })
        except Exception:
            continue
    return entries


def _extract_text_result(state: MainState, role_name: str) -> str:
    try:
        result = state.agent_results.get(role_name, {}).get("results", {})
        if isinstance(result, dict):
            return result.get("text") or result.get("raw") or result.get("content") or ""
        if isinstance(result, str):
            return result
    except Exception:
        return ""
    return ""


def _try_parse_json(raw: str) -> Dict[str, Any] | None:
    if not raw:
        return None
    try:
        return json.loads(raw)
    except Exception:
        pass
    try:
        start = raw.find("{")
        end = raw.rfind("}")
        if start >= 0 and end > start:
            return json.loads(raw[start:end + 1])
    except Exception:
        return None
    return None


def _style_hint(style: str) -> str:
    return {
        "insight": "以洞察/要点为主，突出关键发现与启发",
        "analysis": "以分析/推理为主，强调因果、结构化论证",
    }.get(style, "洞察与分析并重")


def _length_hint(length: str) -> str:
    return {
        "short": "简短精炼（约 5-8 段）",
        "standard": "标准长度（约 8-12 段）",
        "long": "详细深入（约 12-18 段）"
    }.get(length, "标准长度")


@register("kb_report")
def create_kb_report_graph() -> GenericGraphBuilder:
    builder = GenericGraphBuilder(state_model=KBReportState, entry_point="_start_")

    def _start_(state: KBReportState) -> KBReportState:
        if not state.request.file_paths:
            state.request.file_paths = []
        if not state.result_path:
            project_root = get_project_root()
            ts = int(time.time())
            email = getattr(state.request, "email", "") or "default"
            output_dir = project_root / "outputs" / "kb_outputs" / email / f"{ts}_report"
            output_dir.mkdir(parents=True, exist_ok=True)
            state.result_path = str(output_dir)
        state.file_entries = []
        state.file_summaries = []
        state.report_outline = ""
        state.report_markdown = ""
        return state

    async def parse_files_node(state: KBReportState) -> KBReportState:
        state.file_entries = _build_file_entries(state.request.file_paths or [])
        return state

    async def summarize_files_node(state: KBReportState) -> KBReportState:
        if not state.file_entries:
            return state

        prompt_prefix = [
            "你是一位文档分析师，请对给定文档内容进行结构化总结。",
            "输出 JSON，格式如下：",
            "{",
            '  "summary": "...",',
            '  "key_points": ["...", "..."],',
            '  "evidence": [{"quote": "...", "section": "..."}]',
            "}",
            f"输出语言：{state.request.language}",
        ]

        async def _summarize(entry: Dict[str, Any], idx: int) -> Dict[str, Any]:
            agent_name = "kb_prompt_agent"
            prompt = "\n".join(
                prompt_prefix
                + [
                    f"\n文档名称：{entry.get('file_name')}",
                    "文档内容（可能截断）：",
                    entry.get("content", "")
                ]
            )
            agent = create_agent(
                name=agent_name,
                model_name=state.request.model,
                chat_api_url=state.request.chat_api_url,
                temperature=0.2,
                parser_type="text",
            )
            temp_state = MainState(request=state.request)
            res_state = await agent.execute(temp_state, prompt=prompt)
            raw = _extract_text_result(res_state, agent_name) or ""
            parsed = _try_parse_json(raw)
            return {
                "file_name": entry.get("file_name"),
                "summary": parsed or {"raw": raw},
                "raw": raw
            }

        tasks = [_summarize(entry, i + 1) for i, entry in enumerate(state.file_entries)]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        summaries: List[Dict[str, Any]] = []
        for item in results:
            if isinstance(item, Exception):
                log.error(f"[kb_report] summarize failed: {item}")
                continue
            summaries.append(item)
        state.file_summaries = summaries
        return state

    async def plan_report_node(state: KBReportState) -> KBReportState:
        if not state.file_summaries:
            return state

        summary_lines = []
        for item in state.file_summaries:
            payload = item.get("raw") or json.dumps(item.get("summary"), ensure_ascii=False)
            summary_lines.append(f"文件：{item.get('file_name')}\n{payload}")

        prompt = f"""你是一位报告规划师，请根据多个文档摘要拟定报告大纲。
要求：
1. 输出清晰的 Markdown 目录结构；
2. 结合多份文档的交叉主题；
3. 报告风格：{_style_hint(state.request.report_style)};
4. 报告长度：{_length_hint(state.request.length)};
5. 输出语言：{state.request.language}

文档摘要：
{chr(10).join(summary_lines)}
"""
        agent = create_agent(
            name="kb_prompt_agent",
            model_name=state.request.model,
            chat_api_url=state.request.chat_api_url,
            temperature=0.2,
            parser_type="text",
        )
        temp_state = MainState(request=state.request)
        res_state = await agent.execute(temp_state, prompt=prompt)
        state.report_outline = _extract_text_result(res_state, "kb_prompt_agent") or ""
        return state

    async def draft_report_node(state: KBReportState) -> KBReportState:
        if not state.file_summaries:
            return state

        summary_lines = []
        for item in state.file_summaries:
            payload = item.get("raw") or json.dumps(item.get("summary"), ensure_ascii=False)
            summary_lines.append(f"文件：{item.get('file_name')}\n{payload}")

        prompt = f"""你是一位专业研究分析师，请根据以下多文档摘要生成一份高质量报告。
要求：
1. 使用 Markdown 输出；
2. 报告风格：{_style_hint(state.request.report_style)};
3. 报告长度：{_length_hint(state.request.length)};
4. 输出语言：{state.request.language};
5. 关键结论后用 [文件名] 形式标注来源；
6. 结构建议：摘要、关键发现、深入分析、结论与建议。

报告大纲（若有）：
{state.report_outline or "无"}

文档摘要：
{chr(10).join(summary_lines)}
"""
        agent = create_agent(
            name="kb_prompt_agent",
            model_name=state.request.model,
            chat_api_url=state.request.chat_api_url,
            temperature=0.2,
            parser_type="text",
        )
        temp_state = MainState(request=state.request)
        res_state = await agent.execute(temp_state, prompt=prompt)
        state.report_markdown = _extract_text_result(res_state, "kb_prompt_agent") or ""
        return state

    async def finalize_report_node(state: KBReportState) -> KBReportState:
        if not state.report_markdown:
            return state

        prompt = f"""请对以下 Markdown 报告进行润色与一致性修订，保持结构与引用格式不变。
要求：
1. 保持 Markdown 结构；
2. 不删改引用标记 [文件名]；
3. 输出语言：{state.request.language}

报告内容：
{state.report_markdown}
"""
        agent = create_agent(
            name="kb_prompt_agent",
            model_name=state.request.model,
            chat_api_url=state.request.chat_api_url,
            temperature=0.1,
            parser_type="text",
        )
        temp_state = MainState(request=state.request)
        res_state = await agent.execute(temp_state, prompt=prompt)
        polished = _extract_text_result(res_state, "kb_prompt_agent") or ""
        if polished:
            state.report_markdown = polished

        try:
            report_path = Path(state.result_path) / "report.md"
            report_path.write_text(state.report_markdown or "", encoding="utf-8")
        except Exception as e:
            log.error(f"Failed to save report: {e}")

        return state

    nodes = {
        "_start_": _start_,
        "parse_files": parse_files_node,
        "summarize_files": summarize_files_node,
        "plan_report": plan_report_node,
        "draft_report": draft_report_node,
        "finalize_report": finalize_report_node,
        "_end_": lambda s: s
    }

    edges = [
        ("_start_", "parse_files"),
        ("parse_files", "summarize_files"),
        ("summarize_files", "plan_report"),
        ("plan_report", "draft_report"),
        ("draft_report", "finalize_report"),
        ("finalize_report", "_end_")
    ]

    builder.add_nodes(nodes).add_edges(edges)
    return builder
