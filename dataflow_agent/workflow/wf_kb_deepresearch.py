from __future__ import annotations

import asyncio
import json
import time
from pathlib import Path
from typing import List, Dict, Any
from urllib.parse import urlparse

import fitz

from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger
from dataflow_agent.state import KBDeepResearchState, MainState
from dataflow_agent.agentroles import create_agent
from dataflow_agent.utils import get_project_root
from dataflow_agent.toolkits.research_tools import search_web, fetch_page_text

log = get_logger(__name__)

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def _safe_domain(url: str) -> str:
    try:
        parsed = urlparse(url)
        return parsed.netloc or url
    except Exception:
        return url


def _resolve_local_path(path_or_url: str) -> Path:
    if not path_or_url:
        return Path()
    if path_or_url.startswith("http"):
        parsed = urlparse(path_or_url)
        path_or_url = parsed.path
    if path_or_url.startswith("/outputs/"):
        rel = path_or_url[len("/outputs/") :].lstrip("/")
        return (get_project_root() / "outputs" / rel).resolve()
    p = Path(path_or_url)
    if p.is_absolute():
        return p.resolve()
    return (get_project_root() / p).resolve()


def _extract_file_text(path: Path) -> str:
    if not path.exists():
        return ""
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            doc = fitz.open(path)
            text = ""
            for page in doc:
                text += page.get_text() + "\n"
            return text
        if suffix in {".docx", ".doc"}:
            if Document is None:
                return ""
            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs])
        if suffix in {".pptx", ".ppt"}:
            if Presentation is None:
                return ""
            prs = Presentation(path)
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        if suffix in {".md", ".txt"}:
            return path.read_text(encoding="utf-8", errors="ignore")
        return ""
    except Exception:
        return ""


def _build_context(file_paths: List[str], max_chars: int = 60000) -> str:
    if not file_paths:
        return ""
    chunks: List[str] = []
    for raw in file_paths:
        try:
            local_path = _resolve_local_path(raw)
            if not local_path.exists():
                continue
            text = _extract_file_text(local_path)
            if not text:
                continue
            chunks.append(f"=== {local_path.name} ===\n{text}")
        except Exception:
            continue
    combined = "\n\n".join(chunks)
    if len(combined) > max_chars:
        combined = combined[:max_chars] + "\n\n[...content truncated]"
    return combined


def _parse_queries(raw: str, max_queries: int) -> List[str]:
    if not raw:
        return []
    raw = raw.strip()
    # Try JSON first
    try:
        data = json.loads(raw)
        if isinstance(data, dict) and isinstance(data.get("queries"), list):
            queries = [str(q).strip() for q in data["queries"] if str(q).strip()]
            return queries[:max_queries]
        if isinstance(data, list):
            queries = [str(q).strip() for q in data if str(q).strip()]
            return queries[:max_queries]
    except Exception:
        pass

    # Fallback: parse lines
    lines = [line.strip("-* \t") for line in raw.splitlines()]
    queries = [line for line in lines if line and len(line) > 3]
    return queries[:max_queries]


def _extract_text_result(state: MainState, role_name: str) -> str:
    try:
        result = state.agent_results.get(role_name, {}).get("results", {})
        if isinstance(result, dict):
            return result.get("text") or result.get("raw") or result.get("content") or ""
        if isinstance(result, str):
            return result
    except Exception:
        return ""
    return ""


@register("kb_deepresearch")
def create_kb_deepresearch_graph() -> GenericGraphBuilder:
    """
    Workflow for Knowledge Base Deep Research
    Steps:
    1. Parse KB files (optional)
    2. Web search via SerpAPI (optional)
    3. Generate Markdown report
    """
    builder = GenericGraphBuilder(state_model=KBDeepResearchState, entry_point="_start_")

    def _start_(state: KBDeepResearchState) -> KBDeepResearchState:
        if not state.request.file_paths:
            state.request.file_paths = []
        if not state.result_path:
            project_root = get_project_root()
            ts = int(time.time())
            email = getattr(state.request, "email", "") or "default"
            output_dir = project_root / "outputs" / "kb_outputs" / email / f"{ts}_deepresearch"
            output_dir.mkdir(parents=True, exist_ok=True)
            state.result_path = str(output_dir)
        state.context_text = ""
        state.sub_reports = []
        state.search_results = []
        state.plan_queries = []
        state.sources = []
        state.page_texts = []
        state.report_markdown = ""
        return state

    async def parse_files_node(state: KBDeepResearchState) -> KBDeepResearchState:
        state.context_text = _build_context(state.request.file_paths or [])
        return state

    async def plan_search_node(state: KBDeepResearchState) -> KBDeepResearchState:
        if not state.request.enable_agentic:
            state.plan_queries = [state.request.topic or "research"]
            return state

        max_queries = max(1, int(state.request.max_queries or 6))
        depth = max(1, int(state.request.search_depth or 2))
        prompt = f"""你是一位搜索规划专家。请根据研究主题生成检索子问题。
要求：
1. 输出 JSON，格式为 {{\"queries\": [\"...\", \"...\"]}}
2. queries 数量不超过 {max_queries}，可以覆盖不同维度
3. 语言：{state.request.language}
4. 如果有文件内容，请结合补充更具体的查询

研究主题：{state.request.topic or '未命名主题'}
检索轮次：{depth}
文件摘要（可能截断）：
{state.context_text[:3000] if state.context_text else "无"}
"""

        agent = create_agent(
            name="kb_prompt_agent",
            model_name=state.request.model,
            chat_api_url=state.request.chat_api_url,
            temperature=0.2,
            parser_type="text",
        )
        temp_state = MainState(request=state.request)
        res_state = await agent.execute(temp_state, prompt=prompt)
        raw = _extract_text_result(res_state, "kb_prompt_agent")
        queries = _parse_queries(raw, max_queries=max_queries)
        if not queries:
            queries = [state.request.topic or "research"]
        state.plan_queries = queries
        return state

    async def model_research_node(state: KBDeepResearchState) -> KBDeepResearchState:
        if state.request.mode != "llm":
            return state

        queries = state.plan_queries or [state.request.topic or "research"]
        max_queries = max(1, int(state.request.max_queries or len(queries)))
        queries = queries[:max_queries]

        base_prompt = [
            "你是一位联网研究助手，具备实时检索能力。",
            "请针对给定子问题进行检索与归纳，输出结构化结果：",
            "1) 关键结论（要点）",
            "2) 关键证据/数据",
            "3) 相关来源链接（URL，尽量提供可点击的来源）",
            f"输出语言：{state.request.language}",
        ]
        if state.context_text:
            base_prompt.append("可参考资料（来自用户文件，可能截断）：")
            base_prompt.append(state.context_text[:3000])

        async def _run_query(idx: int, query: str) -> Dict[str, Any]:
            agent_name = f"kb_research_agent_{idx}"
            prompt = "\n".join(base_prompt + [f"\n子问题：{query}\n"])
            agent = create_agent(
                name=agent_name,
                model_name=state.request.model,
                chat_api_url=state.request.chat_api_url,
                temperature=0.2,
                parser_type="text",
            )
            temp_state = MainState(request=state.request)
            res_state = await agent.execute(temp_state, prompt=prompt)
            content = _extract_text_result(res_state, agent_name) or ""
            return {"query": query, "content": content}

        tasks = [_run_query(i + 1, q) for i, q in enumerate(queries)]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        sub_reports: List[Dict[str, Any]] = []
        for idx, item in enumerate(results):
            if isinstance(item, Exception):
                log.error(f"[kb_deepresearch] model_research failed: {item}")
                continue
            if item and isinstance(item, dict):
                sub_reports.append(item)

        state.sub_reports = sub_reports
        return state

    async def web_search_node(state: KBDeepResearchState) -> KBDeepResearchState:
        if state.request.mode != "web":
            return state
        if not state.request.search_api_key:
            state.search_results = []
            return state

        queries = state.plan_queries or [state.request.topic or "research"]
        top_k = max(1, int(state.request.top_k_per_query or 5))
        merged: List[Dict[str, Any]] = []
        summaries: List[Dict[str, Any]] = []
        seen = set()
        for q in queries:
            try:
                payload = await search_web(
                    provider=state.request.search_provider,
                    query=q,
                    api_key=state.request.search_api_key,
                    engine=state.request.search_engine,
                    num=max(top_k, int(state.request.search_num or top_k)),
                    google_cse_id=getattr(state.request, "google_cse_id", "") or None,
                    brave_enable_summarizer=getattr(state.request, "brave_summarizer", False),
                )
            except Exception as e:
                log.error(f"[kb_deepresearch] search_web failed: {e}")
                continue

            results = payload.get("results", []) or []
            if payload.get("summary"):
                summaries.append({
                    "query": q,
                    "summary": payload.get("summary")
                })

            for item in results[:top_k]:
                url = item.get("url") or ""
                key = url.lower().strip() or item.get("title", "").lower().strip()
                if key in seen:
                    continue
                seen.add(key)
                merged.append(item)

        state.search_results = merged
        state.summaries = summaries
        return state

    async def fetch_pages_node(state: KBDeepResearchState) -> KBDeepResearchState:
        if state.request.mode != "web":
            return state
        top_n = max(1, int(state.request.fetch_top_n or 8))
        max_chars = max(1000, int(state.request.max_page_chars or 8000))
        sources: List[Dict[str, Any]] = []
        page_texts: List[Dict[str, Any]] = []
        for idx, item in enumerate(state.search_results[:top_n], start=1):
            url = item.get("url") or ""
            text = await fetch_page_text(url, max_chars=max_chars)
            sources.append({
                "index": idx,
                "title": item.get("title") or "",
                "url": url,
                "snippet": item.get("snippet") or "",
                "source": item.get("source") or _safe_domain(url),
                "extracted_text_len": len(text or "")
            })
            if text:
                page_texts.append({
                    "index": idx,
                    "url": url,
                    "title": item.get("title") or "",
                    "content": text
                })

        state.sources = sources
        state.page_texts = page_texts
        return state

    async def generate_report_node(state: KBDeepResearchState) -> KBDeepResearchState:
        topic = state.request.topic or "未命名主题"
        prompt_parts = [
            "你是一位高级研究分析师，请根据给定主题生成结构化的深度研究报告。",
            "要求：",
            "1. 使用 Markdown 输出，包含清晰的标题和要点。",
            "2. 输出结构建议包含：概览、关键概念、核心观点、证据/数据、争议/不足、结论与建议。",
            f"3. 输出语言为: {state.request.language}",
        ]

        if state.request.mode == "web" and (state.search_results or state.page_texts):
            prompt_parts.append("4. 基于以下检索结果与页面正文生成报告，并在关键结论后用 [n] 标注引用。")
        if state.request.mode == "llm" and state.sub_reports:
            prompt_parts.append("4. 基于以下子问题研究结果进行综合归纳，并尽量保留来源链接。")

        prompt_parts.append(f"\n研究主题：{topic}\n")

        if state.context_text:
            prompt_parts.append("参考资料（来自已选文件，可能被截断）：\n" + state.context_text)

        if state.request.mode == "web" and state.search_results:
            sr_lines = []
            for idx, item in enumerate(state.search_results, start=1):
                sr_lines.append(f"[{idx}] {item.get('title')} - {item.get('url')}\n{item.get('snippet','')}")
            prompt_parts.append("检索结果：\n" + "\n\n".join(sr_lines))

        if state.request.mode == "web" and state.page_texts:
            page_lines = []
            for item in state.page_texts:
                page_lines.append(f"[{item.get('index')}] {item.get('title')} - {item.get('url')}\n{item.get('content')}")
            prompt_parts.append("页面正文（已截断）：\n" + "\n\n".join(page_lines))

        if state.request.mode == "web" and state.summaries:
            summary_lines = []
            for item in state.summaries:
                summary_lines.append(f"Query: {item.get('query')}\nSummary: {item.get('summary')}")
            prompt_parts.append("Brave Summarizer 汇总：\n" + "\n\n".join(summary_lines))

        if state.request.mode == "llm" and state.sub_reports:
            sub_lines = []
            for item in state.sub_reports:
                sub_lines.append(f"子问题：{item.get('query')}\n结果：\n{item.get('content')}")
            prompt_parts.append("子问题研究结果：\n" + "\n\n".join(sub_lines))

        prompt = "\n".join(prompt_parts)

        agent = create_agent(
            name="kb_prompt_agent",
            model_name=state.request.model,
            chat_api_url=state.request.chat_api_url,
            temperature=0.3,
            parser_type="text",
        )
        temp_state = MainState(request=state.request)
        res_state = await agent.execute(temp_state, prompt=prompt)
        state.report_markdown = _extract_text_result(res_state, "kb_prompt_agent") or ""

        try:
            report_path = Path(state.result_path) / "report.md"
            report_path.write_text(state.report_markdown or "", encoding="utf-8")
        except Exception as e:
            log.error(f"Failed to save report: {e}")

        return state

    nodes = {
        "_start_": _start_,
        "parse_files": parse_files_node,
        "plan_search": plan_search_node,
        "model_research": model_research_node,
        "web_search": web_search_node,
        "fetch_pages": fetch_pages_node,
        "generate_report": generate_report_node,
        "_end_": lambda s: s
    }

    edges = [
        ("_start_", "parse_files"),
        ("parse_files", "plan_search"),
        ("plan_search", "model_research"),
        ("model_research", "web_search"),
        ("web_search", "fetch_pages"),
        ("fetch_pages", "generate_report"),
        ("generate_report", "_end_")
    ]

    builder.add_nodes(nodes).add_edges(edges)
    return builder
