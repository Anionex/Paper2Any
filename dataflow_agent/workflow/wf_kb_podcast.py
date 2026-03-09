from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any

import fitz  # PyMuPDF
from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger
from dataflow_agent.state import KBPodcastState, MainState
from dataflow_agent.agentroles import create_agent
from dataflow_agent.utils import get_project_root
import re
import wave
from dataflow_agent.toolkits.multimodaltool.req_tts import (
    generate_speech_bytes_async,
    split_tts_text,
    split_tts_sentences,
    split_tts_text_by_bytes
)

log = get_logger(__name__)

# Try importing office libraries
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

@register("kb_podcast")
def create_kb_podcast_graph() -> GenericGraphBuilder:
    """
    Workflow for Knowledge Base Podcast Generation
    Steps:
    1. Parse uploaded files (PDF/Office)
    2. Generate podcast script using LLM
    3. Generate audio using TTS
    """
    builder = GenericGraphBuilder(state_model=KBPodcastState, entry_point="_start_")

    def _extract_text_result(state: MainState, role_name: str) -> str:
        try:
            result = state.agent_results.get(role_name, {}).get("results", {})
            if isinstance(result, dict):
                return result.get("text") or result.get("raw") or ""
            if isinstance(result, str):
                return result
        except Exception:
            return ""
        return ""

    def _start_(state: KBPodcastState) -> KBPodcastState:
        # Ensure request fields
        if not state.request.files:
            state.request.files = []

        # Initialize output directory
        if not state.result_path:
            project_root = get_project_root()
            import time
            ts = int(time.time())
            email = getattr(state.request, 'email', 'default')
            output_dir = project_root / "outputs" / "kb_outputs" / email / f"{ts}_podcast"
            output_dir.mkdir(parents=True, exist_ok=True)
            state.result_path = str(output_dir)
        else:
            Path(state.result_path).mkdir(parents=True, exist_ok=True)

        state.file_contents = []
        state.podcast_script = ""
        state.audio_path = ""
        return state

    async def parse_files_node(state: KBPodcastState) -> KBPodcastState:
        """
        Parse all files and extract content
        """
        files = state.request.files
        if not files:
            state.file_contents = []
            return state

        async def process_file(file_path: str) -> Dict[str, Any]:
            file_path_obj = Path(file_path)
            filename = file_path_obj.name

            if not file_path_obj.exists():
                return {
                    "filename": filename,
                    "content": f"[Error: File not found {file_path}]"
                }

            suffix = file_path_obj.suffix.lower()
            raw_content = ""

            try:
                # PDF
                if suffix == ".pdf":
                    try:
                        doc = fitz.open(file_path)
                        text = ""
                        for page in doc:
                            text += page.get_text() + "\n"
                        raw_content = text
                    except Exception as e:
                        raw_content = f"[Error parsing PDF: {e}]"

                # Word
                elif suffix in [".docx", ".doc"]:
                    if Document is None:
                         raw_content = "[Error: python-docx not installed]"
                    else:
                        try:
                            doc = Document(file_path)
                            raw_content = "\n".join([p.text for p in doc.paragraphs])
                        except Exception as e:
                             raw_content = f"[Error parsing Docx: {e}]"

                # PPT
                elif suffix in [".pptx", ".ppt"]:
                    if Presentation is None:
                        raw_content = "[Error: python-pptx not installed]"
                    else:
                        try:
                            prs = Presentation(file_path)
                            text = ""
                            for i, slide in enumerate(prs.slides):
                                text += f"--- Slide {i+1} ---\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text += shape.text + "\n"
                            raw_content = text
                        except Exception as e:
                            raw_content = f"[Error parsing PPT: {e}]"

                else:
                    try:
                        with open(file_path, "r", encoding="utf-8") as f:
                            raw_content = f.read()
                    except:
                        raw_content = "[Unsupported file type]"

            except Exception as e:
                 raw_content = f"[Parse Error: {e}]"

            # Truncate content
            truncated_content = raw_content[:50000] if len(raw_content) > 50000 else raw_content

            return {
                "filename": filename,
                "content": truncated_content
            }

        # Run in parallel
        tasks = [process_file(f) for f in files]
        results = await asyncio.gather(*tasks)

        state.file_contents = results
        return state

    async def generate_script_node(state: KBPodcastState) -> KBPodcastState:
        """
        Generate podcast script using LLM
        """
        if not state.file_contents:
            state.podcast_script = "No content available for podcast generation."
            return state

        # Format file contents
        contents_str = ""
        for item in state.file_contents:
            contents_str += f"=== {item['filename']} ===\n{item['content']}\n\n"

        # Podcast script prompt
        language = state.request.language
        mode = getattr(state.request, "podcast_mode", "monologue")
        length_pref = str(getattr(state.request, "podcast_length", "standard")).lower()
        length_alias = {
            "brief": "brief",
            "short": "brief",
            "精炼": "brief",
            "精简": "brief",
            "简短": "brief",
            "standard": "standard",
            "normal": "standard",
            "默认": "standard",
            "long": "long",
            "详细": "long",
            "扩展": "long",
            "加长": "long",
        }
        length_pref = length_alias.get(length_pref, "standard")
        if length_pref == "brief":
            length_hint_zh = "时长约2-4分钟，内容精炼，删去铺垫与重复。"
            length_hint_en = "Target 2–4 minutes, concise and no padding."
        elif length_pref == "long":
            length_hint_zh = "时长约10-15分钟，内容更完整，可加入适量例子与细节。"
            length_hint_en = "Target 10–15 minutes, more detail with examples."
        else:
            length_hint_zh = "时长约5-10分钟，内容完整但不过度铺垫。"
            length_hint_en = "Target 5–10 minutes, complete but not verbose."
        if mode == "dialog":
            speaker_a = "主持人" if language == "zh" else "Host"
            speaker_b = "嘉宾" if language == "zh" else "Guest"
            prompt = f"""你是一位专业的知识播客制作人。基于以下资料，生成一段双人对话播客脚本。

要求：
1. 口语化、生动有趣，避免书面语
2. 结构清晰：开场白 → 核心内容 → 总结
3. 使用类比和例子帮助理解
4. 适当加入互动性语言（"你可能会想..."）
5. 使用{language}语言
6. {length_hint_zh if language == "zh" else length_hint_en}
7. 严格使用如下格式逐行输出（每行一个角色）：
{speaker_a}: ...
{speaker_b}: ...

资料内容：
{contents_str}

请生成播客脚本：
输出要求：（speaker后面不要有任何md格式内容，直接纯文本）
"""
        else:
            prompt = f"""你是一位专业的知识播客主播。基于以下资料，生成一段知识播客脚本。

要求：
1. 口语化、生动有趣，避免书面语
2. 结构清晰：开场白 → 核心内容 → 总结
3. 使用类比和例子帮助理解
4. 适当加入互动性语言（"你可能会想..."）
5. 使用{language}语言
6. {length_hint_zh if language == "zh" else length_hint_en}

资料内容：
{contents_str}

请生成播客脚本：
输出要求：（不要有任何md格式内容，直接纯文本返回！）
"""

        try:
            agent = create_agent(
                name="kb_prompt_agent",
                model_name=state.request.model,
                chat_api_url=state.request.chat_api_url,
                temperature=0.7,
                parser_type="text"
            )

            temp_state = MainState(request=state.request)
            res_state = await agent.execute(temp_state, prompt=prompt)

            state.podcast_script = _extract_text_result(res_state, "kb_prompt_agent") or "[Script generation failed]"
        except Exception as e:
            log.error(f"Script generation failed: {e}")
            state.podcast_script = f"[Script generation error: {e}]"

        # Save script to file
        try:
            script_path = Path(state.result_path) / "script.txt"
            script_path.write_text(state.podcast_script, encoding="utf-8")
        except Exception as e:
            log.error(f"Failed to save script: {e}")

        return state

    async def generate_audio_node(state: KBPodcastState) -> KBPodcastState:
        """
        Generate audio using TTS
        """
        if not state.podcast_script or state.podcast_script.startswith("["):
            state.audio_path = ""
            return state

        try:
            audio_path = str(Path(state.result_path) / "podcast.wav")
            mode = getattr(state.request, "podcast_mode", "monologue")
            tts_model = getattr(state.request, "tts_model", "")
            tts_model_l = str(tts_model).lower()
            is_openai_tts = tts_model_l.startswith("gpt-4o-mini-tts") or tts_model_l.startswith("tts-1")
            is_cosyvoice_tts = "cosyvoice" in tts_model_l

            max_chars: int | None = None
            max_bytes: int | None = None
            if is_openai_tts:
                max_chars = 3800  # OpenAI TTS input limit is 4096 chars
            elif is_cosyvoice_tts:
                max_bytes = 4000  # CosyVoice: reduced byte limit to avoid truncation
            else:
                max_chars = 3000

            if mode == "dialog":
                if max_chars:
                    max_chars = min(max_chars, 1200)
                if max_bytes:
                    max_bytes = min(max_bytes, 4000)
                concurrency = 1
            else:
                concurrency = 2
            if is_cosyvoice_tts:
                concurrency = 1

            segments = []
            if mode == "dialog":
                language = state.request.language
                speaker_a = "主持人" if language == "zh" else "Host"
                speaker_b = "嘉宾" if language == "zh" else "Guest"
                speaker_map = {
                    speaker_a.lower(): "A",
                    speaker_b.lower(): "B",
                    "a": "A",
                    "b": "B",
                    "speaker a": "A",
                    "speaker b": "B",
                    "角色a": "A",
                    "角色b": "B",
                    "主播": "A",
                    "嘉宾": "B",
                }
                pattern = re.compile(r"^\s*([^:：]{1,20})\s*[:：]\s*(.+)$")
                current_speaker = "A"
                for raw_line in state.podcast_script.splitlines():
                    line = raw_line.strip()
                    if not line:
                        continue
                    m = pattern.match(line)
                    if m:
                        label = m.group(1).strip().lower()
                        content = m.group(2).strip()
                        mapped = speaker_map.get(label)
                        if mapped:
                            current_speaker = mapped
                        if content:
                            segments.append({"speaker": current_speaker, "text": content})
                        continue
                    # No label, append to current speaker
                    if segments and segments[-1]["speaker"] == current_speaker:
                        segments[-1]["text"] = f"{segments[-1]['text']} {line}"
                    else:
                        segments.append({"speaker": current_speaker, "text": line})

                expanded = []
                for seg in segments:
                    sentences = split_tts_sentences(seg["text"])
                    if not sentences:
                        continue
                    if max_bytes:
                        for chunk in split_tts_text_by_bytes(" ".join(sentences), max_bytes):
                            expanded.append({
                                "speaker": seg["speaker"],
                                "text": chunk
                            })
                    else:
                        for sent in sentences:
                            if max_chars and len(sent) <= max_chars:
                                expanded.append({
                                    "speaker": seg["speaker"],
                                    "text": sent
                                })
                            else:
                                for chunk in split_tts_text(sent, max_chars or 1500):
                                    expanded.append({
                                        "speaker": seg["speaker"],
                                        "text": chunk
                                    })
                segments = expanded
            else:
                if max_bytes:
                    for chunk in split_tts_text_by_bytes(state.podcast_script, max_bytes):
                        segments.append({"speaker": "A", "text": chunk})
                else:
                    for chunk in split_tts_text(state.podcast_script, max_chars or 1500):
                        segments.append({"speaker": "A", "text": chunk})

            if not segments:
                raise RuntimeError("No valid TTS segments generated from script")

            sem = asyncio.Semaphore(concurrency)

            async def _run(seg):
                voice = state.request.voice_name if seg["speaker"] == "A" else state.request.voice_name_b
                async with sem:
                    return await generate_speech_bytes_async(
                        text=seg["text"],
                        api_url=state.request.chat_api_url,
                        api_key=state.request.api_key,
                        model=state.request.tts_model,
                        voice_name=voice,
                    )

            async def _run_no_sem(seg):
                voice = state.request.voice_name if seg["speaker"] == "A" else state.request.voice_name_b
                return await generate_speech_bytes_async(
                    text=seg["text"],
                    api_url=state.request.chat_api_url,
                    api_key=state.request.api_key,
                    model=state.request.tts_model,
                    voice_name=voice,
                )

            async def _run_with_retry(seg, attempts=3, base_delay=0.8, use_sem=True):
                last_err = None
                for i in range(attempts):
                    try:
                        if use_sem:
                            return await _run(seg)
                        return await _run_no_sem(seg)
                    except Exception as e:
                        last_err = e
                        if i < attempts - 1:
                            await asyncio.sleep(base_delay * (i + 1))
                        continue
                raise last_err

            if concurrency <= 1:
                results = []
                for seg in segments:
                    results.append(await _run_with_retry(seg, attempts=2, use_sem=False))
            else:
                tasks = [asyncio.create_task(_run_with_retry(seg, attempts=2, use_sem=True)) for seg in segments]
                results = await asyncio.gather(*tasks, return_exceptions=True)

            failed_indices = [i for i, r in enumerate(results) if isinstance(r, Exception)]
            if failed_indices:
                log.warning(f"TTS retry sequentially for {len(failed_indices)} failed segment(s)")
                for i in failed_indices:
                    results[i] = await _run_with_retry(segments[i], attempts=3, use_sem=False)

            audio_chunks = results

            os.makedirs(os.path.dirname(os.path.abspath(audio_path)), exist_ok=True)
            with wave.open(audio_path, "wb") as wav_file:
                wav_file.setnchannels(1)        # 1 Channel
                wav_file.setsampwidth(2)        # 16 bit = 2 bytes
                wav_file.setframerate(24000)    # 24kHz
                wav_file.writeframes(b"".join(audio_chunks))

            state.audio_path = audio_path
            log.info(f"Audio generated successfully: {audio_path}")
        except Exception as e:
            log.error(f"Audio generation failed: {e}")
            state.audio_path = f"[Audio generation error: {e}]"

        return state

    nodes = {
        "_start_": _start_,
        "parse_files": parse_files_node,
        "generate_script": generate_script_node,
        "generate_audio": generate_audio_node,
        "_end_": lambda s: s
    }

    edges = [
        ("_start_", "parse_files"),
        ("parse_files", "generate_script"),
        ("generate_script", "generate_audio"),
        ("generate_audio", "_end_")
    ]

    builder.add_nodes(nodes).add_edges(edges)
    return builder
